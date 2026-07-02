from __future__ import annotations

import json
import os
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from srxy import SearchResult


FIXTURES_DIR = Path(__file__).resolve().parent / "fixtures"
CORPUS_DIR = FIXTURES_DIR / "corpus"
FILE_SEARCH_DIR = FIXTURES_DIR / "file_search"
OCR_FIXTURES_DIR = FILE_SEARCH_DIR / "ocr"
OCR_IMAGE_FIXTURE = OCR_FIXTURES_DIR / "ocr_sample.png"
OCR_PDF_FIXTURE = OCR_FIXTURES_DIR / "ocr_embedded.pdf"
MINIMAL_JPEG_FIXTURE = FILE_SEARCH_DIR / "minimal.jpg"
MINIMAL_MP3_FIXTURE = FILE_SEARCH_DIR / "minimal.mp3"
MINIMAL_MP4_FIXTURE = FILE_SEARCH_DIR / "minimal.mp4"
PHOTOSHOP_XMP_JPEG_FIXTURE = FILE_SEARCH_DIR / "samples" / "images" / "photoshop_xmp.jpg"


def file_search_root() -> Path:
	override = os.environ.get("SRXY_FILE_SEARCH_FIXTURES", "").strip()
	if override:
		return Path(override).expanduser()
	return FILE_SEARCH_DIR


def file_search_docs() -> Path:
	return file_search_root()


def file_search_samples() -> Path:
	return file_search_root() / "samples"


def require_file_search_fixtures() -> Path:
	root = file_search_root()
	marker = root / "notes.txt"
	if not marker.is_file():
		raise FileNotFoundError(
			f"File-search fixtures not found at {root}. Expected tests/fixtures/file_search/ in the repository checkout."
		)
	return root


def cuda_available() -> bool:
	try:
		import importlib.util

		if importlib.util.find_spec("torch") is None:
			return False
		import torch  # type: ignore[reportMissingImports]

		return bool(torch.cuda.is_available())
	except (ImportError, OSError, RuntimeError):
		return False


def integration_test_cpu_requested(config: object) -> bool:
	getoption = getattr(config, "getoption", None)
	if callable(getoption) and getoption("--integration-test-cpu"):
		return True
	value = os.environ.get("SRXY_INTEGRATION_TEST_CPU", "").strip().lower()
	return value in {"1", "true", "yes", "on"}


def transcribe_device_matrix_devices(config: object) -> list[str]:
	if cuda_available():
		if integration_test_cpu_requested(config):
			return ["cuda", "cpu"]
		return ["cuda"]
	return ["cpu"]


def copy_media_fixture(name: str, destination: Path) -> None:
	destination.write_bytes((FILE_SEARCH_DIR / name).read_bytes())


@dataclass(frozen=True)
class LabeledQuery:
	query: str
	expected_id: str
	min_top_score: float


@dataclass
class Product:
	name: str
	description: str = ""
	category: str = ""
	status: str = ""
	alias: str = ""
	sku: str = ""
	tags: str = ""


def load_search_corpus() -> list[dict[str, Any]]:
	path = CORPUS_DIR / "search_corpus.json"
	return json.loads(path.read_text(encoding="utf-8"))


def load_labeled_queries() -> list[LabeledQuery]:
	path = CORPUS_DIR / "labeled_queries.json"
	raw = json.loads(path.read_text(encoding="utf-8"))
	return [LabeledQuery(**entry) for entry in raw]


def top_k_hit_rate(results_by_query: list[list[SearchResult]], expected_ids: list[str], *, k: int) -> float:
	hits = 0
	for results, expected_id in zip(results_by_query, expected_ids, strict=True):
		top_ids = [result.item["id"] for result in results[:k]]
		if expected_id in top_ids:
			hits += 1
	return hits / len(expected_ids)


def write_pdf_with_text(path: Path, text: str):
	"""Write a minimal single-page PDF with extractable text."""
	escaped = text.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")
	stream = f"BT /F1 12 Tf 72 720 Td ({escaped}) Tj ET"
	stream_bytes = stream.encode("latin-1")
	pdf = f"""%PDF-1.4
1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj
2 0 obj<</Type/Pages/Kids[3 0 R]/Count 1>>endobj
3 0 obj<</Type/Page/Parent 2 0 R/MediaBox[0 0 612 792]/Resources<</Font<</F1 4 0 R>>>>/Contents 5 0 R>>endobj
4 0 obj<</Type/Font/Subtype/Type1/BaseFont/Helvetica>>endobj
5 0 obj<</Length {len(stream_bytes)}>>stream
{stream}
endstream
endobj
xref
0 6
0000000000 65535 f 
0000000009 00000 n 
0000000058 00000 n 
0000000115 00000 n 
0000000244 00000 n 
0000000313 00000 n 
trailer<</Size 6/Root 1 0 R>>
startxref
403
%%EOF
"""
	path.write_bytes(pdf.encode("latin-1"))


def write_docx_with_text(path: Path, text: str):
	from docx import Document

	document = Document()
	document.add_paragraph(text)
	document.save(str(path))


def write_xlsx_with_text(path: Path, text: str):
	from openpyxl import Workbook

	workbook = Workbook()
	sheet = workbook.active
	if sheet is not None:
		sheet.title = "Summary"
		sheet["A1"] = text
	workbook.save(path)


def write_pptx_with_text(path: Path, text: str):
	from pptx import Presentation
	from pptx.util import Inches

	presentation = Presentation()
	slide = presentation.slides.add_slide(presentation.slide_layouts[6])
	textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
	textbox.text_frame.text = text
	presentation.save(str(path))


def _large_embed_image_bytes(image_path: Path) -> bytes:
	import io

	from PIL import Image

	from srxy.ocr_text import MIN_PDF_IMAGE_OCR_BYTES

	source = image_path.read_bytes()
	if len(source) >= MIN_PDF_IMAGE_OCR_BYTES:
		return source

	with Image.open(image_path) as image:
		rgb = image.convert("RGB")
		for scale in (2, 4, 8, 16):
			width = max(rgb.width * scale, 400)
			height = max(rgb.height * scale, 300)
			resized = rgb.resize((width, height))
			buffer = io.BytesIO()
			resized.save(buffer, format="PNG")
			data = buffer.getvalue()
			if len(data) >= MIN_PDF_IMAGE_OCR_BYTES:
				return data
	raise RuntimeError(f"unable to create embedded image >= {MIN_PDF_IMAGE_OCR_BYTES} bytes from {image_path}")


def write_docx_with_image(path: Path, image_path: Path, text: str = ""):
	import io

	from docx import Document

	embed_bytes = _large_embed_image_bytes(image_path)
	document = Document()
	if text:
		document.add_paragraph(text)
	document.add_picture(io.BytesIO(embed_bytes))
	document.save(str(path))


def write_pptx_with_image(path: Path, image_path: Path, text: str = ""):
	import io

	from pptx import Presentation
	from pptx.util import Inches

	embed_bytes = _large_embed_image_bytes(image_path)
	presentation = Presentation()
	slide = presentation.slides.add_slide(presentation.slide_layouts[6])
	if text:
		textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
		textbox.text_frame.text = text
	slide.shapes.add_picture(io.BytesIO(embed_bytes), Inches(1), Inches(2))
	presentation.save(str(path))


def write_xlsx_with_image(path: Path, image_path: Path, text: str = ""):
	import tempfile

	from openpyxl import Workbook
	from openpyxl.drawing.image import Image as XLImage

	embed_bytes = _large_embed_image_bytes(image_path)
	workbook = Workbook()
	sheet = workbook.active
	with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as handle:
		handle.write(embed_bytes)
		temp_image = Path(handle.name)
	try:
		if sheet is not None:
			sheet.title = "Summary"
			if text:
				sheet["A1"] = text
			sheet.add_image(XLImage(str(temp_image)), "A2")
		workbook.save(path)
	finally:
		temp_image.unlink(missing_ok=True)


def write_mp4_with_tags(path: Path, *, title: str | None = None, min_size: int = 0):
	from mutagen.mp4 import MP4

	copy_media_fixture("minimal.mp4", path)
	if title is not None:
		mp4 = MP4(path)
		mp4["\xa9nam"] = [title]
		mp4.save()
	if min_size > 0:
		current_size = path.stat().st_size
		if current_size < min_size:
			with path.open("ab") as handle:
				handle.write(b"\x00" * (min_size - current_size))


def set_windows_tags(path: Path, tags: list[str]) -> None:
	from srxy.windows_metadata import write_windows_keywords

	write_windows_keywords(path, tags)
