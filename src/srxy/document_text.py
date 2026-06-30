from __future__ import annotations

import json
import re
import sys
import xml.etree.ElementTree as ET
from collections.abc import Callable, Iterator
from pathlib import Path
from zipfile import ZipFile

from srxy.archive_guard import ArchiveGuardError, validate_zip_archive


DOCUMENT_SUFFIXES = frozenset({".pdf", ".docx", ".xlsx", ".pptx"})
_OFFICE_SUFFIXES = frozenset({".docx", ".xlsx", ".pptx"})
_DC_NS = "http://purl.org/dc/elements/1.1/"
_CP_NS = "http://schemas.openxmlformats.org/package/2006/metadata/core-properties"
_APP_NS = "http://schemas.openxmlformats.org/officeDocument/2006/extended-properties"
_CORE_METADATA_FIELDS = {
	f"{{{_DC_NS}}}creator": "Author",
	f"{{{_CP_NS}}}lastModifiedBy": "Last saved by",
	f"{{{_DC_NS}}}title": "Title",
	f"{{{_DC_NS}}}subject": "Subject",
	f"{{{_CP_NS}}}keywords": "Keywords",
	f"{{{_DC_NS}}}description": "Description",
	f"{{{_CP_NS}}}category": "Category",
	f"{{{_CP_NS}}}contentStatus": "Content status",
}
_APP_METADATA_FIELDS = {
	f"{{{_APP_NS}}}Application": "Program name",
	f"{{{_APP_NS}}}Company": "Company",
	f"{{{_APP_NS}}}Manager": "Manager",
}


def is_document_path(path: Path) -> bool:
	return path.suffix.lower() in DOCUMENT_SUFFIXES


def iter_document_metadata_lines(path: Path) -> Iterator[tuple[int, str]]:
	suffix = path.suffix.lower()
	if suffix not in _OFFICE_SUFFIXES:
		return
	if sys.platform == "win32":
		from srxy.windows_metadata import windows_metadata_supported

		if windows_metadata_supported():
			return
	try:
		validate_zip_archive(path)
		entries = _collect_office_package_metadata(path)
	except (ArchiveGuardError, OSError):
		return
	except Exception:
		return
	for line_number, (label, value) in enumerate(entries, start=1):
		yield line_number, f"[{label}] {value}"


def _collect_office_package_metadata(path: Path) -> list[tuple[str, str]]:
	entries: list[tuple[str, str]] = []
	with ZipFile(path) as archive:
		names = set(archive.namelist())
		if "docProps/core.xml" in names:
			entries.extend(_parse_metadata_entries(archive.read("docProps/core.xml"), _CORE_METADATA_FIELDS))
		if "docProps/app.xml" in names:
			entries.extend(_parse_metadata_entries(archive.read("docProps/app.xml"), _APP_METADATA_FIELDS))
	return entries


def _parse_metadata_entries(raw: bytes, field_labels: dict[str, str]) -> list[tuple[str, str]]:
	try:
		root = ET.fromstring(raw)  # noqa: S314
	except ET.ParseError:
		return []

	entries: list[tuple[str, str]] = []
	for tag, label in field_labels.items():
		element = root.find(f".//{tag}")
		if element is None or element.text is None:
			continue
		text = element.text.strip()
		if not text:
			continue
		if label == "Keywords":
			for part in re.split(r"[;,]", text):
				part = part.strip()
				if part:
					entries.append((label, part))
			continue
		entries.append((label, text))
	return entries


def _document_cache_variant(path: Path, *, ocr: bool | None = None) -> str:
	from srxy.ocr_text import is_ocr_active

	suffix = path.suffix.lower()
	if suffix == ".pdf":
		return f"{suffix}:ocr={int(is_ocr_active(ocr))}"
	return suffix


def _encode_document_lines(lines: list[tuple[int, str, str]]) -> bytes:
	return "\n".join(
		json.dumps([line_number, text, location_kind]) for line_number, text, location_kind in lines
	).encode("utf-8")


def _decode_document_lines(payload: bytes) -> list[tuple[int, str, str]]:
	lines: list[tuple[int, str, str]] = []
	for raw_line in payload.decode("utf-8").splitlines():
		if not raw_line:
			continue
		line_number, text, location_kind = json.loads(raw_line)
		lines.append((int(line_number), str(text), str(location_kind)))
	return lines


def iter_document_lines(path: Path, *, ocr: bool | None = None) -> Iterator[tuple[int, str, str]]:
	suffix = path.suffix.lower()
	extractors: dict[str, Callable[..., Iterator[tuple[int, str, str]]]] = {
		".pdf": _iter_pdf_lines,
		".docx": _iter_docx_lines,
		".xlsx": _iter_xlsx_lines,
		".pptx": _iter_pptx_lines,
	}
	extractor = extractors.get(suffix)
	if extractor is None:
		return

	from srxy.cache import CACHE_KIND_DOCUMENT_TEXT, cache_get, cache_put, get_file_content_hash

	try:
		content_hash = get_file_content_hash(path)
		variant = _document_cache_variant(path, ocr=ocr)
		cached = cache_get(CACHE_KIND_DOCUMENT_TEXT, content_hash, variant)
		if cached is not None:
			yield from _decode_document_lines(cached)
			return

		if suffix == ".pdf":
			lines = list(extractor(path, ocr=ocr))
		else:
			lines = list(extractor(path))
		cache_put(CACHE_KIND_DOCUMENT_TEXT, content_hash, variant, _encode_document_lines(lines))
		yield from lines
	except ArchiveGuardError:
		return
	except Exception:
		return


def _iter_pdf_lines(path: Path, *, ocr: bool | None = None) -> Iterator[tuple[int, str, str]]:
	from pypdf import PdfReader

	from srxy.ocr_text import is_ocr_active, ocr_max_file_size, ocr_pdf_page_images

	reader = PdfReader(path)
	ocr_active = is_ocr_active(ocr)
	if ocr_active:
		try:
			limit = ocr_max_file_size()
			if limit is not None and path.stat().st_size > limit:
				ocr_active = False
		except OSError:
			ocr_active = False

	for page_number, page in enumerate(reader.pages, start=1):
		embedded = (page.extract_text() or "").strip()
		image_ocr = ocr_pdf_page_images(page).strip() if ocr_active else ""
		if embedded:
			yield page_number, embedded, "page"
		if image_ocr:
			yield page_number, image_ocr, "ocr"


def _iter_docx_lines(path: Path) -> Iterator[tuple[int, str, str]]:
	from docx import Document

	validate_zip_archive(path)
	document = Document(str(path))
	for paragraph_number, paragraph in enumerate(document.paragraphs, start=1):
		text = paragraph.text.strip()
		if text:
			yield paragraph_number, text, "paragraph"


def _iter_xlsx_lines(path: Path) -> Iterator[tuple[int, str, str]]:
	from openpyxl import load_workbook

	validate_zip_archive(path)
	workbook = load_workbook(path, read_only=True, data_only=True)
	try:
		line_number = 0
		for sheet in workbook.worksheets:
			for row in sheet.iter_rows(values_only=True):
				cells = [str(cell) for cell in row if cell is not None and str(cell).strip()]
				if not cells:
					continue
				line_number += 1
				yield line_number, f"[{sheet.title}] " + " ".join(cells), "row"
	finally:
		workbook.close()


def _iter_pptx_lines(path: Path) -> Iterator[tuple[int, str, str]]:
	from pptx import Presentation

	validate_zip_archive(path)
	presentation = Presentation(str(path))
	for slide_number, slide in enumerate(presentation.slides, start=1):
		parts: list[str] = []
		for shape in slide.shapes:
			text = shape.text.strip() if hasattr(shape, "text") else ""
			if text:
				parts.append(text)
		if parts:
			yield slide_number, " ".join(parts), "slide"
